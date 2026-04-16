"""
Betonarme Döşeme Tasarım Sistemi - Sunum Oluşturucu
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Renk Paleti ──────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0A, 0x1B, 0x3D)   # koyu lacivert (arka plan)
BLUE_MID  = RGBColor(0x1A, 0x50, 0x8B)   # orta mavi
ACCENT    = RGBColor(0x00, 0xAE, 0xEF)   # parlak cyan (vurgu)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRY = RGBColor(0xE8, 0xF1, 0xFB)
GOLD      = RGBColor(0xFF, 0xC0, 0x00)
GREEN     = RGBColor(0x00, 0xC8, 0x6E)
ORANGE    = RGBColor(0xFF, 0x6B, 0x35)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]   # tamamen boş düzen

# ── Yardımcı Fonksiyonlar ────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb, alpha=None, line_rgb=None, line_width=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_textbox_multiline(slide, lines, l, t, w, h,
                          size=14, bold=False, color=WHITE,
                          align=PP_ALIGN.LEFT, line_spacing_pt=6):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(line_spacing_pt)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txb

def slide_background(slide, color=NAVY):
    """Slayt arka planını tek renk yap."""
    add_rect(slide, 0, 0, 13.33, 7.5, color)

def header_bar(slide, title, subtitle=None):
    """Üst şerit + başlık."""
    add_rect(slide, 0, 0, 13.33, 1.3, BLUE_MID)
    add_rect(slide, 0, 1.3, 13.33, 0.06, ACCENT)
    add_text(slide, title, 0.4, 0.08, 12.5, 0.75,
             size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.78, 12.5, 0.45,
                 size=16, color=LIGHT_GRY, align=PP_ALIGN.LEFT)

def footer_bar(slide, text="Betonarme Döşeme Otomatik Hesap Sistemi  |  TS500"):
    add_rect(slide, 0, 7.1, 13.33, 0.4, BLUE_MID)
    add_text(slide, text, 0.3, 7.1, 12.5, 0.4,
             size=10, color=LIGHT_GRY, align=PP_ALIGN.CENTER)

def icon_card(slide, icon, title, body_lines, l, t, w=3.8, h=2.4,
              card_color=BLUE_MID, accent=ACCENT):
    """Simge + başlık + madde içerikli kart."""
    add_rect(slide, l, t, w, h, card_color, line_rgb=accent, line_width=1.5)
    # Simge dairesi
    circ = slide.shapes.add_shape(9, Inches(l+0.15), Inches(t+0.15),
                                   Inches(0.55), Inches(0.55))
    circ.fill.solid(); circ.fill.fore_color.rgb = accent
    circ.line.fill.background()
    tf2 = circ.text_frame
    p2  = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2  = p2.add_run(); r2.text = icon
    r2.font.size = Pt(18); r2.font.bold = True; r2.font.color.rgb = NAVY

    add_text(slide, title, l+0.8, t+0.18, w-1.0, 0.45,
             size=14, bold=True, color=WHITE)
    add_textbox_multiline(slide, body_lines, l+0.2, t+0.75, w-0.4,
                          h-0.9, size=12, color=LIGHT_GRY)

def numbered_step(slide, number, title, desc, l, t, w=11.8, h=0.72,
                  bg=BLUE_MID, num_bg=ACCENT):
    """Numaralı adım satırı."""
    add_rect(slide, l, t, w, h, bg)
    nb = slide.shapes.add_shape(9, Inches(l+0.08), Inches(t+0.08),
                                  Inches(0.55), Inches(0.55))
    nb.fill.solid(); nb.fill.fore_color.rgb = num_bg; nb.line.fill.background()
    tf = nb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(number)
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = NAVY

    add_text(slide, title, l+0.75, t+0.06, 3.5, 0.35,
             size=14, bold=True, color=WHITE)
    add_text(slide, desc, l+0.75, t+0.37, w-1.0, 0.3,
             size=11, color=LIGHT_GRY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLAYT 1 – KAPAK
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK_LAYOUT)
slide_background(s1, NAVY)

# Sol dekoratif şerit
add_rect(s1, 0, 0, 0.35, 7.5, ACCENT)
add_rect(s1, 0.35, 0, 0.1, 7.5, BLUE_MID)

# Sağ dekoratif kutu
add_rect(s1, 9.5, 1.0, 3.5, 5.5, BLUE_MID)
add_rect(s1, 9.5, 1.0, 3.5, 0.06, ACCENT)
add_rect(s1, 9.5, 6.44, 3.5, 0.06, ACCENT)

# Başlık
add_text(s1, "YAPay ZEKA DESTEKLİ", 0.7, 1.1, 8.5, 0.8,
         size=22, bold=False, color=ACCENT)
add_text(s1, "Betonarme Döşeme\nOtomatik Hesap Sistemi", 0.7, 1.75, 8.5, 2.0,
         size=40, bold=True, color=WHITE)
add_text(s1, "YOLO Nesne Tanıma  ·  OCR Metin Okuma  ·  TS500 Mühendislik Hesabı  ·  DXF/PDF Çıktı",
         0.7, 3.65, 8.5, 0.7, size=15, color=LIGHT_GRY)

# Sağ kutu içeriği
add_text(s1, "🏗", 9.8, 1.4, 2.8, 1.2, size=60, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox_multiline(s1,
    ["📌 Kalıp Planı → Hesap",
     "🤖 YOLOv8 Segmentasyon",
     "📐 TS500 / Koef. Yöntemi",
     "📄 DXF + PDF Çıktı"],
    9.65, 2.7, 3.2, 2.8, size=13, color=WHITE)

add_rect(s1, 0.7, 5.6, 8.5, 0.06, ACCENT)
add_text(s1, "Bitirme Projesi Sunumu  ·  2025–2026", 0.7, 5.7, 8.5, 0.5,
         size=14, color=LIGHT_GRY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLAYT 2 – GENEL BAKIŞ / PROBLEMİN TANIMI
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK_LAYOUT)
slide_background(s2)
header_bar(s2, "Projenin Amacı ve Motivasyonu",
           "Kalıp planlarından donatı hesabına kadar tam otomatik bir akış")
footer_bar(s2)

# Sol – Problem
add_rect(s2, 0.3, 1.6, 5.9, 5.2, BLUE_MID)
add_rect(s2, 0.3, 1.6, 5.9, 0.06, ORANGE)
add_text(s2, "⚠️  Mevcut Sorun", 0.55, 1.7, 5.5, 0.5, size=16, bold=True, color=ORANGE)
add_textbox_multiline(s2, [
    "• Mühendisler kalıp planlarını manuel olarak inceliyor",
    "• Her döşeme için ayrı ayrı момент ve donatı hesabı yapılıyor",
    "• Bu süreç zaman alıcı ve insan hatalarına açık",
    "• DXF/AutoCAD çizimi de ek çaba gerektiriyor",
    "",
    "Büyük projelerde onlarca döşeme binlerce değişken",
    "içerebilmektedir. Manuel süreç verimli değildir.",
], 0.55, 2.3, 5.5, 4.2, size=13, color=LIGHT_GRY)

# Sağ – Çözüm
add_rect(s2, 6.9, 1.6, 5.9, 5.2, BLUE_MID)
add_rect(s2, 6.9, 1.6, 5.9, 0.06, GREEN)
add_text(s2, "✅  Geliştirilen Çözüm", 7.15, 1.7, 5.5, 0.5, size=16, bold=True, color=GREEN)
add_textbox_multiline(s2, [
    "• Kalıp planı görseli sisteme yüklenir",
    "• YOLOv8 modeli döşeme, kiriş, balkon,",
    "  aks balonu gibi sınıfları otomatik tanır",
    "• EasyOCR aks ve boyut metinlerini okur",
    "• TS500'e göre moment ve donatı hesaplanır",
    "• Donatı planı DXF + PDF olarak çıktılanır",
    "",
    "→ Tüm akış tek tıkla, saniyeler içinde!",
], 7.15, 2.3, 5.5, 4.2, size=13, color=LIGHT_GRY)

add_rect(s2, 6.5, 1.6, 0.08, 5.2, ACCENT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLAYT 3 – SİSTEM MİMARİSİ (GENEL AKIŞDİYAGRAMI)
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK_LAYOUT)
slide_background(s3)
header_bar(s3, "Sistem Mimarisi", "Uçtan uca pipeline – 5 ana katman")
footer_bar(s3)

stages = [
    ("📷", "GİRDİ",       "Kalıp planı\ngörseli (PNG/JPG)",    ACCENT),
    ("🤖", "YOLOv8",      "Nesne tespiti\n9 sınıf segmentasyon", GOLD),
    ("🔤", "OCR",         "EasyOCR ile\naks & boyut okuma",     GREEN),
    ("📐", "HESAP",       "TS500 moment\n& donatı tasarımı",    ORANGE),
    ("📄", "ÇIKTI",       "DXF donatı planı\n+ PDF raporu",     ACCENT),
]

# Ok şeritler
for i, (icon, title, body, clr) in enumerate(stages):
    lx = 0.35 + i * 2.55
    # Kutu
    add_rect(s3, lx, 1.6, 2.35, 4.5, BLUE_MID, line_rgb=clr, line_width=2)
    add_rect(s3, lx, 1.6, 2.35, 0.08, clr)
    # İkon dairesi
    c = s3.shapes.add_shape(9, Inches(lx+0.85), Inches(1.85),
                              Inches(0.65), Inches(0.65))
    c.fill.solid(); c.fill.fore_color.rgb = clr; c.line.fill.background()
    tf = c.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = icon; r.font.size = Pt(20)
    r.font.color.rgb = NAVY if clr != NAVY else WHITE
    # Başlık
    add_text(s3, title, lx+0.1, 2.6, 2.15, 0.45,
             size=15, bold=True, color=clr, align=PP_ALIGN.CENTER)
    # İçerik
    add_text(s3, body, lx+0.1, 3.15, 2.15, 2.8,
             size=12, color=LIGHT_GRY, align=PP_ALIGN.CENTER)
    # Ok (son kutu için yok)
    if i < len(stages) - 1:
        arr = s3.shapes.add_shape(13,   # rightArrow
                                   Inches(lx+2.35), Inches(3.6),
                                   Inches(0.2), Inches(0.3))
        arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT; arr.line.fill.background()

# Alt açıklama
add_rect(s3, 0.35, 6.3, 12.6, 0.55, BLUE_MID)
add_text(s3,
    "Her adım birbirinin çıktısını girdi olarak kullanır – insan müdahalesi gerektirmez.",
    0.55, 6.32, 12.2, 0.5, size=13, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLAYT 4 – YOLO MODELİ
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK_LAYOUT)
slide_background(s4)
header_bar(s4, "Adım 1 – YOLOv8 Nesne Tespiti",
           "Kalıp planındaki yapısal elemanların otomatik tanınması")
footer_bar(s4)

# Sol: sınıf listesi
add_rect(s4, 0.3, 1.55, 5.8, 5.25, BLUE_MID)
add_text(s4, "Tespit Edilen 9 Sınıf", 0.55, 1.62, 5.3, 0.45,
         size=15, bold=True, color=GOLD)

classes = [
    ("0", "slab_area",       "Döşeme alanı"),
    ("1", "beam_area",       "Kiriş alanı"),
    ("2", "balcony_area",    "Balkon alanı"),
    ("3", "axis_bubble",     "Aks balonu"),
    ("4", "dimension_text",  "Boyut metni"),
    ("5", "panel_text",      "Panel etiketi (D11...)"),
    ("6", "balcony_text",    "Balkon etiketi (BL...)"),
    ("7", "column_symbol",   "Kolon sembolü"),
    ("8", "opening_area",    "Boşluk / açıklık"),
]
clr_map = {0:GREEN, 1:ORANGE, 2:ACCENT, 3:GOLD,
           4:LIGHT_GRY, 5:LIGHT_GRY, 6:LIGHT_GRY,
           7:ORANGE, 8:RGBColor(0xFF,0x50,0x50)}

for i, (idx, name, desc) in enumerate(classes):
    ty = 2.15 + i * 0.52
    clr = clr_map.get(int(idx), WHITE)
    add_rect(s4, 0.45, ty, 0.35, 0.38, clr)
    add_text(s4, idx, 0.45, ty, 0.35, 0.38,
             size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s4, f"{name}", 0.9, ty+0.02, 2.0, 0.35, size=11, bold=True, color=clr)
    add_text(s4, desc, 2.95, ty+0.02, 3.0, 0.35, size=10, color=LIGHT_GRY)

# Sağ: model detayları
add_rect(s4, 6.8, 1.55, 6.2, 2.5, BLUE_MID)
add_text(s4, "Model Detayları", 7.05, 1.62, 5.8, 0.45,
         size=15, bold=True, color=ACCENT)
add_textbox_multiline(s4, [
    "🏋  Model Ağırlıkları   :  best.pt  (~38 MB)",
    "📐  Giriş Çözünürlüğü  :  1024 × 1024 piksel",
    "🎯  Güven Eşiği        :  conf = 0.35 (ayarlanabilir)",
    "🔧  Kütüphane          :  Ultralytics YOLOv8",
    "📊  Görev              :  Object Detection / Segmentation",
], 7.05, 2.15, 5.8, 1.8, size=13, color=LIGHT_GRY)

add_rect(s4, 6.8, 4.2, 6.2, 2.6, BLUE_MID)
add_text(s4, "Nasıl Çalışır?", 7.05, 4.27, 5.8, 0.45,
         size=15, bold=True, color=GREEN)
add_textbox_multiline(s4, [
    "1. Görsel piksel düzeyinde analiz edilir",
    "2. Her tespit için sınıf + güven skoru + bbox üretilir",
    "3. Bounding box koordinatları akslar arası hücre",
    "   ataması için kullanılır",
    "4. Bakma/karıştırma: slab ↔ balcony ↔ opening",
    "   birbirine karıştırılmadan ayrılır",
], 7.05, 4.82, 5.8, 1.85, size=13, color=LIGHT_GRY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLAYT 5 – AKS ÇIKARIMI & OCR
# ═══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK_LAYOUT)
slide_background(s5)
header_bar(s5, "Adım 2 – Aks Çıkarımı & OCR Metin Okuma",
           "Yapısal grid'in ve boyutların otomatik olarak belirlenmesi")
footer_bar(s5)

add_textbox_multiline(s5, [
    "Hibrit Aks Oluşturma (build_axes_hybrid):",
    "   1️⃣  Kolon sembol merkezleri toplanır",
    "   2️⃣  Döşeme sınır koordinatları eklenir",
    "   3️⃣  Aks balonu merkezleri eklenir",
    "   4️⃣  Yakın koordinatlar birleştirilir (tol = 30 px)",
    "   5️⃣  OCR etiketleri en yakın aks çizgisine atanır",
    "   6️⃣  Eksik etiketler otomatik doldurulur (1,2,3… / A,B,C…)",
], 0.4, 1.6, 6.3, 3.6, size=13.5, color=WHITE)

add_textbox_multiline(s5, [
    "EasyOCR – Boyut Okuma (parse_spans_cm_from_ocr):",
    "   • Sayısal tokenlar (20–9999 cm arası) ayıklanır",
    "   • Küresel piksel/cm ölçeği tahmin edilir",
    "   • Her aks aralığına en uyumlu değer atanır",
    "   • Fiziksel tutarlılık denetimi (±%50 tolerans)",
    "   • Hatalı okumalar (kiriş etiketi, kesir) filtrelenir",
], 0.4, 5.2, 6.3, 2.1, size=13.5, color=LIGHT_GRY)

# Sağ: hücre tipi atama
add_rect(s5, 6.9, 1.6, 6.1, 5.7, BLUE_MID)
add_text(s5, "Hücre Tipi Atama", 7.1, 1.65, 5.8, 0.45,
         size=15, bold=True, color=ACCENT)
add_textbox_multiline(s5, [
    "Her (i, j) hücre için:",
    "",
    "  📗  SLAB    → Hücrenin %15'inden fazlası",
    "               slab_area tespiti ile örtüşüyor",
    "",
    "  📕  OPENING  → opening_area tespiti var",
    "",
    "  ⬜  VOID     → Hiçbiri yok (duvar / dış alan)",
    "",
    "Bu grid mantığıyla her döşeme  bbox_ij  olarak",
    "kodlanır: [i0, j0, i1, j1]",
    "",
    "Span'lar cm → grid index'e çevrilir:",
    "   count = round(span_cm / dx_cm)",
    "   edges = kümülatif toplu liste",
], 7.1, 2.2, 5.8, 5.0, size=12.5, color=LIGHT_GRY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLAYT 6 – OCR SINIRLAMASI: KUCUK DOSEMELER
# ═══════════════════════════════════════════════════════════════════════════════
s6_ocr = prs.slides.add_slide(BLANK_LAYOUT)
slide_background(s6_ocr)
header_bar(s6_ocr, "OCR Sınırlaması – Küçük Döşeme Metin Tespiti",
           "Büyük döşemede iyi çalışırken küçük döşeme metinleri neden zor okunuyor?")
footer_bar(s6_ocr)

# ---------- Sol üst blok: SEBEP -----------------------------------------------
add_rect(s6_ocr, 0.3, 1.58, 6.1, 2.55, BLUE_MID, line_rgb=ORANGE, line_width=2)
add_rect(s6_ocr, 0.3, 1.58, 6.1, 0.07, ORANGE)
add_text(s6_ocr, "SORUN   –   Neden Başarısız?", 0.55, 1.65, 5.7, 0.42,
         size=15, bold=True, color=ORANGE)
add_textbox_multiline(s6_ocr, [
    "Büyük döşeme (örn. 8x6 m):",
    "  Boyut metni ekranda GENİŞ alana yayılır",
    "  → Piksel yoğunluğu yüksek → OCR rahat okur",
    "",
    "Küçük döşeme (örn. 2x1.5 m):",
    "  Aynı metin çok dar bir bant içine sıkışır",
    "  → Piksel yoğunluğu düşük → harfler bulanık",
    "  → EasyOCR karakterleri birleştirip yanlış okur",
], 0.55, 2.15, 5.7, 1.9, size=12.5, color=LIGHT_GRY)

# ---------- Sol alt blok: MEVCUT ÖNLEM ----------------------------------------
add_rect(s6_ocr, 0.3, 4.3, 6.1, 2.5, BLUE_MID, line_rgb=GOLD, line_width=2)
add_rect(s6_ocr, 0.3, 4.3, 6.1, 0.07, GOLD)
add_text(s6_ocr, "MEVCUT ÖNLEM   (Sistemde Var Olan)", 0.55, 4.37, 5.7, 0.42,
         size=14, bold=True, color=GOLD)
add_textbox_multiline(s6_ocr, [
    "_preprocess_variants() fonksiyonu:",
    "  V1  Orijinal kırpıntı x3 büyütme (INTER_CUBIC)",
    "  V2  CLAHE: Yerel kontrast artırma + büyütme",
    "  V3  Otomatik eşikleme → görüntüyü siyah/beyaza",
    "       çevirme (koyu metin, açık zemin ayrışır)",
    "  V4  Gürültü temizleme → ince lekeler silinir,",
    "       metin pikselleri netleştirilir",
    "  Oylama: En çok onaylanan değer kazanır",
], 0.55, 4.85, 5.7, 1.85, size=12.5, color=LIGHT_GRY)

# ---------- Sağ: ÇÖZÜM ÖNERİLERİ --------------------------------------------
add_rect(s6_ocr, 6.75, 1.58, 6.25, 5.22, BLUE_MID, line_rgb=ACCENT, line_width=2)
add_rect(s6_ocr, 6.75, 1.58, 6.25, 0.07, ACCENT)
add_text(s6_ocr, "ÇÖZÜM ÖNERİLERİ", 7.0, 1.65, 5.8, 0.42,
         size=15, bold=True, color=ACCENT)

solutions = [
    ("1", GREEN,
     "Giriş Çözünürlüğünü Artır",
     "imgsz = 1024 → 2048 piksel ile yeniden çalıştır;",
     "küçük nesneler daha fazla piksel alanı kazanır"),
    ("2", GOLD,
     "Adaptif Kırpıntı Büyütme",
     "Bbox alanı < eşik ise büyütme katsayısı",
     "3x → 6x'e çıkarılır (bbox boyutuna göre dinamik)"),
    ("3", ORANGE,
     "Yapay Zeka ile Süper Çözünürlük",
     "EDSR / real-ESRGAN modeli ile düşük çözünürlüklü",
     "kırpıntıyı 4x büyütme → metin pikselleri artar"),
    ("4", GREEN,
     "GUI Manuel Düzeltme",
     "Tespit edilemeyen boyut değeri kullanıcı tarafından",
     "JSON / GUI üzerinden elle girilebilir"),
    ("5", ACCENT,
     "Boyut Metnine Özel YOLO Modeli",
     "Yalnızca boyut metnine odaklı ayrı bir",
     "YOLO modeli eğitilerek hassasiyet artırılır"),
]

for i, (num, clr, title, line1, line2) in enumerate(solutions):
    ty = 2.18 + i * 0.93
    nb = s6_ocr.shapes.add_shape(9, Inches(6.9), Inches(ty+0.1),
                                   Inches(0.5), Inches(0.5))
    nb.fill.solid(); nb.fill.fore_color.rgb = clr; nb.line.fill.background()
    tf = nb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(16)
    r.font.bold = True; r.font.color.rgb = NAVY
    add_text(s6_ocr, title, 7.5, ty+0.06, 5.3, 0.35,
             size=12.5, bold=True, color=clr)
    add_text(s6_ocr, line1, 7.5, ty+0.4, 5.3, 0.28,
             size=11, color=LIGHT_GRY)
    add_text(s6_ocr, line2, 7.5, ty+0.65, 5.3, 0.25,
             size=11, color=LIGHT_GRY)

# Alt bilgi çubuğu
add_rect(s6_ocr, 0.3, 6.95, 12.7, 0.36, RGBColor(0x00, 0x2B, 0x55))
add_text(s6_ocr,
    "Not: 8x6 m döşemede çalışmasının temel nedeni → geniş alan = yüksek piksel yoğunluğu = OCR doğruluğu artar",
    0.5, 6.96, 12.3, 0.34, size=11, italic=True,
    color=GOLD, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLAYT 7 – TS500 HESAP
# ═══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK_LAYOUT)
slide_background(s6)
header_bar(s6, "Adım 3 – TS500 Mühendislik Hesabı",
           "Moment, donatı seçimi ve mesnet dengelemesi")
footer_bar(s6)

# Döşeme tipleri kart satırı
for i, (icon, title, lines, clr) in enumerate([
    ("↔", "TEK DOĞRULTULU\n(ONEWAY)", [
        "• Ll / Ls > 2 koşulu",
        "• Ana + dağıtım donatısı",
        "• Konsol (mesnet) momenti",
        "  TS500 Katsayı Yöntemi",
        "• Pilye ayrıştırması",
    ], ACCENT),
    ("⊞", "ÇİFT DOĞRULTULU\n(TWOWAY)", [
        "• Ll / Ls ≤ 2 koşulu",
        "• Mx ve My doğrultularında",
        "  ayrı moment hesabı",
        "• Mesnet dengelemesi",
        "  (balance_support_moments)",
        "• 4 kenar için As seçimi",
    ], GOLD),
    ("↗", "KONSOL / BALKON\n(BALCONY)", [
        "• Ankastre mesnetli",
        "• Negatif moment",
        "  M = q·L²/2  formülü",
        "• Üst donatı tasarımı",
        "• TS500 min. donatı",
        "  kontrolleri",
    ], GREEN),
]):
    lx = 0.3 + i * 4.35
    add_rect(s6, lx, 1.6, 4.1, 5.2, BLUE_MID, line_rgb=clr, line_width=2)
    add_rect(s6, lx, 1.6, 4.1, 0.08, clr)
    add_text(s6, icon, lx+1.55, 1.8, 1.0, 0.8,
             size=32, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_text(s6, title, lx+0.15, 2.65, 3.8, 0.7,
             size=14, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_textbox_multiline(s6, lines, lx+0.2, 3.4, 3.7, 3.2,
                          size=12.5, color=LIGHT_GRY)

# Alt bar: ortak parametreler
add_rect(s6, 0.3, 6.85, 12.7, 0.45, BLUE_MID)
add_text(s6,
    "Ortak Girdiler: Beton sınıfı (C20–C50)  ·  Çelik sınıfı (B420C / B500C)  ·  h (döşeme kalınlığı mm)  ·  cover (pas payı mm)  ·  Pd (kN/m²)",
    0.5, 6.86, 12.3, 0.42, size=11, color=LIGHT_GRY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLAYT 7 – DOĞRULTUDA ÇALIŞAN DÖŞEME VE MESNET DENGELEMESİ
# ═══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK_LAYOUT)
slide_background(s7)
header_bar(s7, "Mesnet Dengelemesi – Komşu Döşeme Etkileşimi",
           "Paylaşılan mesnetlerde farklı moment değerlerinin dengelenmesi (TS500)")
footer_bar(s7)

add_rect(s7, 0.3, 1.6, 12.7, 5.2, BLUE_MID)

add_text(s7, "Problem:", 0.6, 1.7, 4.0, 0.45, size=15, bold=True, color=ORANGE)
add_textbox_multiline(s7, [
    "Komşu iki döşeme aynı mesnet kirişini paylaşıyor.",
    "Her biri kendi başına hesaplanırken farklı negatif",
    "moment değerleri üretir → çelişkili donatı!",
], 0.6, 2.2, 5.5, 1.6, size=13, color=LIGHT_GRY)

add_text(s7, "TS500 Çözümü:", 0.6, 3.85, 4.0, 0.45,
         size=15, bold=True, color=GREEN)
add_textbox_multiline(s7, [
    "1. Her iki döşemenin mesnet momenti hesaplanır",
    "2. Büyük olan değer %30 azaltılır (Mmax × 0,70)",
    "3. Küçük olan değer %30 artırılır (Mmin × 1,30)",
    "4. İkisi de eşit yapılır → tek donatı seçimi",
    "5. Bu işlem tüm ortak mesnet çiftleri için yapılır",
], 0.6, 4.35, 5.7, 2.2, size=13, color=LIGHT_GRY)

# Sağ şema
add_rect(s7, 7.0, 1.7, 5.8, 5.0, NAVY)
add_text(s7, "Dengeleme Şeması", 7.2, 1.75, 5.4, 0.45,
         size=14, bold=True, color=ACCENT)

# Döşeme A
add_rect(s7, 7.3, 2.3, 2.1, 1.2, BLUE_MID, line_rgb=ACCENT, line_width=1.5)
add_text(s7, "Döşeme A", 7.3, 2.3, 2.1, 0.5,
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s7, "M_neg = 8.5 kNm", 7.3, 2.8, 2.1, 0.5,
         size=11, color=ORANGE, align=PP_ALIGN.CENTER)

# Ortak kiriş
add_rect(s7, 9.45, 2.1, 0.2, 1.6, ACCENT)
add_text(s7, "KİRİŞ", 9.4, 2.1, 0.35, 1.6,
         size=9, color=NAVY, align=PP_ALIGN.CENTER)

# Döşeme B
add_rect(s7, 9.7, 2.3, 2.1, 1.2, BLUE_MID, line_rgb=GOLD, line_width=1.5)
add_text(s7, "Döşeme B", 9.7, 2.3, 2.1, 0.5,
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s7, "M_neg = 6.2 kNm", 9.7, 2.8, 2.1, 0.5,
         size=11, color=GOLD, align=PP_ALIGN.CENTER)

add_text(s7, "↓ Dengeleme", 8.5, 3.7, 2.0, 0.45,
         size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

add_rect(s7, 7.3, 4.2, 4.5, 0.7, RGBColor(0x00,0x40,0x20))
add_text(s7, "M_denge = 8.5×0.70 = 5.95  ≈  6.2×1.30 = 8.06  →  Ortak: ~7.0 kNm",
         7.4, 4.25, 4.3, 0.6, size=11, color=GREEN, align=PP_ALIGN.CENTER)

add_text(s7, "Her iki döşeme aynı mesnet donatısı kullanır → tutarlılık sağlanır",
         7.3, 5.0, 5.5, 0.55, size=12, italic=True, color=LIGHT_GRY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLAYT 8 – DXF & PDF ÇIKTISI
# ═══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK_LAYOUT)
slide_background(s8)
header_bar(s8, "Adım 4 – DXF Donatı Planı & PDF Çıktısı",
           "AutoCAD uyumlu çizim dosyası ve otomatik PDF dönüşümü")
footer_bar(s8)

add_textbox_multiline(s8, [
    "📐  DXF Çıktısı (dxf_out.py)",
    "",
    "• Her döşeme için ayrı çizim katmanı (Layer)",
    "• Duz + pilye donatıları ayrı renk/tip çizgi",
    "• Aks hatları ve boyut çizgileri",
    "• Kiriş konumları kalın çizgi",
    "• Donatı isimleri, aralıkları metin olarak",
    "• Balkon donatısı üst katta gösterilir",
    "",
    "📄  Otomatik PDF (ezdxf + matplotlib)",
    "",
    "• ezdxf kütüphanesi DXF'i okur",
    "• MatplotlibBackend A4 yatay sayfaya render eder",
    "• Koyu arka plan → tüm çizgiler net görünür",
    "• Kullanıcı tek tıkla PDF'i açabilir",
], 0.4, 1.6, 6.1, 5.7, size=13, color=LIGHT_GRY)

# Sağ: layer listesi
add_rect(s8, 7.0, 1.6, 5.9, 5.7, BLUE_MID)
add_text(s8, "DXF Katman Yapısı", 7.2, 1.67, 5.5, 0.45,
         size=15, bold=True, color=ACCENT)

layers = [
    (GREEN,  "SLAB_EDGE",      "Döşeme kenar çizgileri"),
    (ACCENT, "AXIS",           "Aks çizgileri"),
    (GOLD,   "BEAM",           "Kiriş çizgileri"),
    (ORANGE, "REBAR_DUZ",      "Düz donatı"),
    (ACCENT, "REBAR_PILYE",    "Pilye (kıvrık) donatı"),
    (GREEN,  "REBAR_DIST",     "Dağıtım donatısı"),
    (LIGHT_GRY,"DIMENSION",    "Boyut ve etiket metinleri"),
    (RGBColor(0xFF,0x50,0x50),"BALCONY_EDGE","Balkon kenrı"),
],

for i, (clr, name, desc) in enumerate(layers[0]):
    ty = 2.22 + i * 0.62
    add_rect(s8, 7.15, ty, 0.18, 0.38, clr)
    add_text(s8, name, 7.4, ty+0.02, 2.6, 0.35, size=11, bold=True, color=clr)
    add_text(s8, desc, 10.05, ty+0.02, 2.7, 0.35, size=10, color=LIGHT_GRY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLAYT 9 – GUI KULLANIM AKIŞI
# ═══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK_LAYOUT)
slide_background(s9)
header_bar(s9, "Kullanıcı Arayüzü – Kullanım Akışı",
           "Tkinter tabanlı GUI – adım adım işlem")
footer_bar(s9)

steps = [
    ("1", "Görsel Yükle (AI)", "Kalıp planı PNG/JPG seçilir → YOLO+OCR pipeline otomatik çalışır",      ACCENT),
    ("2", "Malzeme Gir",       "Beton sınıfı, çelik sınıfı, h (mm), pas payı seçilir",                  GOLD),
    ("3", "Kirişleri Düzenle", "Otomatik tespite ek olarak kullanıcı kiriş ekleyip çıkarabilir",         GREEN),
    ("4", "Hesapla",           "'Hesapla' butonuna basılır → TS500 hesabı yapılır, rapor ekranda görünür", ORANGE),
    ("5", "DXF/PDF Al",        "Donatı planı DXF olarak kaydedilir, PDF otomatik oluşturulur",            ACCENT),
]

for i, (num, title, desc, clr) in enumerate(steps):
    ty = 1.62 + i * 1.0
    add_rect(s9, 0.3, ty, 12.5, 0.85, BLUE_MID)
    add_rect(s9, 0.3, ty, 12.5, 0.06, clr)
    nb = s9.shapes.add_shape(9, Inches(0.42), Inches(ty+0.12),
                               Inches(0.6), Inches(0.6))
    nb.fill.solid(); nb.fill.fore_color.rgb = clr; nb.line.fill.background()
    tf = nb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(20)
    r.font.bold = True; r.font.color.rgb = NAVY

    add_text(s9, title, 1.15, ty+0.1, 3.2, 0.4,
             size=14, bold=True, color=clr)
    add_text(s9, desc, 4.4, ty+0.1, 8.1, 0.6,
             size=13, color=LIGHT_GRY)

# Alt not
add_rect(s9, 0.3, 6.9, 12.5, 0.4, NAVY)
add_text(s9,
    "💡  'JSON'dan Yükle' butonu ile önceden kaydedilmiş veri de kullanılabilir  ·  "
    "Conf ayarı ile tespit hassasiyeti anlık değiştirilebilir",
    0.5, 6.91, 12.1, 0.38, size=11, color=LIGHT_GRY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLAYT 10 – MODÜL HARİTASI
# ═══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK_LAYOUT)
slide_background(s10)
header_bar(s10, "Kod Mimarisi – Modül Haritası",
           "Projenin dosya yapısı ve sorumluluk dağılımı")
footer_bar(s10)

modules = [
    ("main.py",                "Giriş noktası",            "Uygulamayı başlatır",                          ACCENT),
    ("gui.py",                 "Kullanıcı Arayüzü",        "Tkinter GUI, tıklama olayları, çizim",          GOLD),
    ("pipeline.py",            "Orkestrasyoncu",           "YOLO → OCR → JSON → Slab akışını yönetir",     GREEN),
    ("infer_to_calc_inputs.py","Yapay Zeka Motoru",        "YOLO çalıştırır, aks çıkarır, OCR okur",        ORANGE),
    ("json_loader.py",         "Veri Ayrıştırıcı",         "JSON → RealSlab + BeamEdge dönüşümü",           ACCENT),
    ("slab_model.py",          "Veri Modeli",              "SlabSystem, Slab sınıfları ve hesap mantığı",   GOLD),
    ("oneway_slab.py",         "Tek Doğrultulu Hesap",     "TS500 koeffisient yöntemi – As tasarımı",       GREEN),
    ("twoway_slab.py",         "Çift Doğrultulu Hesap",    "Mx / My momentleri ve donatı seçimi",           ORANGE),
    ("balcony_slab.py",        "Balkon Hesabı",            "Konsol moment ve üst donatı tasarımı",          ACCENT),
    ("moment_balance_slab.py", "Mesnet Dengelemesi",       "Komşu döşeme mesnet momentlerini dengeler",     GOLD),
    ("struct_design.py",       "Mühendislik Formülleri",   "As_min, smax, rebar seçim fonksiyonları",       GREEN),
    ("dxf_out.py",             "DXF Yazıcı",               "AutoCAD uyumlu çizim dosyası üretir",           ORANGE),
    ("constants.py",           "Sabitler",                 "Beton/çelik tabloları, TS500 katsayıları",      ACCENT),
]

# 2 sütun
col1 = modules[:7]
col2 = modules[7:]

for i, (fname, role, desc, clr) in enumerate(col1):
    ty = 1.68 + i * 0.75
    add_rect(s10, 0.3, ty, 6.1, 0.65, BLUE_MID)
    add_rect(s10, 0.3, ty, 0.12, 0.65, clr)
    add_text(s10, fname, 0.52, ty+0.04, 2.4, 0.3, size=11, bold=True, color=clr)
    add_text(s10, role, 0.52, ty+0.34, 2.4, 0.28, size=9, color=GOLD)
    add_text(s10, desc, 2.95, ty+0.1, 3.3, 0.45, size=10, color=LIGHT_GRY)

for i, (fname, role, desc, clr) in enumerate(col2):
    ty = 1.68 + i * 0.75
    add_rect(s10, 6.9, ty, 6.1, 0.65, BLUE_MID)
    add_rect(s10, 6.9, ty, 0.12, 0.65, clr)
    add_text(s10, fname, 7.12, ty+0.04, 2.55, 0.3, size=11, bold=True, color=clr)
    add_text(s10, role, 7.12, ty+0.34, 2.55, 0.28, size=9, color=GOLD)
    add_text(s10, desc, 9.7, ty+0.1, 3.1, 0.45, size=10, color=LIGHT_GRY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLAYT 11 – ÖZET / SONUÇ
# ═══════════════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK_LAYOUT)
slide_background(s11, NAVY)
add_rect(s11, 0, 0, 0.35, 7.5, ACCENT)

add_text(s11, "SONUÇ & KATKILAR", 0.7, 0.8, 11.5, 0.7,
         size=28, bold=True, color=ACCENT)
add_rect(s11, 0.7, 1.5, 11.5, 0.06, ACCENT)

highlights = [
    ("🤖", "Yapay Zeka Entegrasyonu",
     "YOLOv8 ile formwork planlarından 9 farklı yapısal eleman sınıfı otomatik tanınır"),
    ("🔤", "Optik Karakter Tanıma",
     "EasyOCR ile aks etiketleri ve boyutlar piksel düzeyinde okunur"),
    ("📐", "TS500 Hesabı",
     "Tek/çift doğrultulu ve konsol döşemeler için tam mühendislik hesabı"),
    ("⚖️", "Mesnet Dengelemesi",
     "Komşu döşemeler arasında TS500'e uygun moment dengeleme algoritması"),
    ("📄", "DXF + PDF Çıktı",
     "AutoCAD uyumlu donatı planı ve otomatik PDF oluşturma"),
    ("⚡", "Tam Otomasyon",
     "Kalıp planından donatı planına kadar tek tıkla, insan müdahalesi minimum"),
]

for i, (icon, title, desc) in enumerate(highlights):
    col = i % 2
    row = i // 2
    lx = 0.7 + col * 6.15
    ty = 1.8 + row * 1.7
    add_rect(s11, lx, ty, 5.85, 1.5, BLUE_MID, line_rgb=ACCENT, line_width=1)
    c = s11.shapes.add_shape(9, Inches(lx+0.12), Inches(ty+0.38),
                               Inches(0.6), Inches(0.6))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.line.fill.background()
    tf = c.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = icon; r.font.size = Pt(18)
    r.font.color.rgb = NAVY
    add_text(s11, title, lx+0.85, ty+0.12, 4.8, 0.45,
             size=14, bold=True, color=WHITE)
    add_text(s11, desc, lx+0.85, ty+0.6, 4.8, 0.8,
             size=11, color=LIGHT_GRY)

footer_bar(s11)

# ── Kaydet ───────────────────────────────────────────────────────────────────
out_path = "Betonarme_Doseme_Sistemi_Sunum.pptx"
prs.save(out_path)
print(f"Sunum kaydedildi: {out_path}")
